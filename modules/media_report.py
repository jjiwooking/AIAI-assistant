import os
import json
import re
import tempfile
import time
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont

MODEL_NAME = "gemini-3.6-flash"

THEMES = {
    "Corporate Basic": {
        "accent": RGBColor(72, 78, 180),
        "dark": RGBColor(30, 35, 48),
        "soft": RGBColor(245, 247, 252),
    },
    "Executive": {
        "accent": RGBColor(31, 78, 121),
        "dark": RGBColor(20, 28, 40),
        "soft": RGBColor(243, 247, 251),
    },
    "Technical Report": {
        "accent": RGBColor(15, 118, 110),
        "dark": RGBColor(29, 42, 44),
        "soft": RGBColor(242, 249, 248),
    },
    "Minimal": {
        "accent": RGBColor(75, 85, 99),
        "dark": RGBColor(31, 41, 55),
        "soft": RGBColor(249, 250, 251),
    },
}


def _get_api_key():
    api_key = os.getenv("GEMINI_API_KEY", "")

    if not api_key:
        try:
            import streamlit as st
            api_key = st.secrets.get("GEMINI_API_KEY", "")
        except Exception:
            api_key = ""

    api_key = str(api_key).strip()

    if not api_key or api_key == "여기에_API_키를_입력하세요":
        raise RuntimeError("GEMINI_API_KEY가 설정되지 않았습니다.")

    return api_key


def _client():
    try:
        from google import genai
    except ImportError as e:
        raise RuntimeError(
            "google-genai가 설치되지 않았습니다. requirements.txt를 확인해주세요."
        ) from e

    return genai.Client(api_key=_get_api_key())


def _extract_json(text):
    if not text:
        raise ValueError("AI 응답이 비어 있습니다.")

    text = str(text).strip()
    text = re.sub(r"^```(?:json)?\s*", "", text)
    text = re.sub(r"\s*```$", "", text)

    try:
        return json.loads(text)

    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, flags=re.S)

        if not match:
            raise ValueError("AI 응답에서 JSON 데이터를 찾지 못했습니다.")

        return json.loads(match.group(0))


def _normalize(data):
    data = data if isinstance(data, dict) else {}

    defaults = {
        "title": "회의·발표 분석",
        "transcript": "",
        "uncertain_terms": [],
        "summary": [],
        "key_points": [],
        "action_items": [],
        "numbers": [],
        "timeline": [],
        "slide_plan": [],
    }

    for key, value in defaults.items():
        data.setdefault(key, value)

    return data


def _media_prompt(analysis_type, language_mode, glossary):
    glossary_text = ", ".join(glossary) if glossary else "없음"

    return f"""
당신은 한국 기업의 회의·발표 분석 전문 AI입니다.
첨부된 음성 또는 영상을 처음부터 끝까지 확인하고 실무 보고용으로 분석하세요.

[분석 유형]
{analysis_type}

[언어]
{language_mode}

[회사/업무 전문용어]
{glossary_text}

[매우 중요한 규칙]
1. 한국어 발화를 자연스러운 한국어 Transcript로 작성합니다.
2. 영어 회사명, 장비명, 약어는 문맥과 전문용어 사전을 우선합니다.
3. 확신이 낮은 표현을 억지로 확정하지 않습니다.
4. 확신이 낮은 표현은 uncertain_terms에 반드시 기록합니다.
5. severity는 yellow 또는 red만 사용합니다.
6. yellow = 확인 권장, red = 오인식 가능성이 높아 수정 권장입니다.
7. uncertain_terms의 text는 transcript 안에 실제 존재하는 문자열과 같아야 합니다.
8. 실제 음성/영상에서 확인된 숫자만 numbers에 넣습니다.
9. 담당자, 기한, 숫자, 결정사항을 임의로 만들어내지 않습니다.
10. 영상의 화면에 실제로 보이는 제목, 표, 수치, 일정도 필요한 경우 반영합니다.
11. 발표 자료는 한국 회사 보고용으로 간결하게 구성합니다.
12. 데이터가 부족하면 시각화를 억지로 만들지 않습니다.

[반환 형식]
아래 형식의 JSON만 반환하세요.

{{
  "title": "분석 제목",
  "transcript": "전체 Transcript",
  "uncertain_terms": [
    {{
      "text": "애매한 원문",
      "suggestion": "추천 표현",
      "severity": "yellow",
      "reason": "확인이 필요한 이유"
    }}
  ],
  "summary": [
    "핵심 요약 1",
    "핵심 요약 2"
  ],
  "key_points": [
    "중요 포인트 1",
    "중요 포인트 2"
  ],
  "action_items": [
    {{
      "owner": "담당 또는 미정",
      "task": "할 일",
      "due": "기한 또는 미정",
      "status": "대기"
    }}
  ],
  "numbers": [
    {{
      "label": "항목명",
      "value": 100,
      "unit": "단위",
      "context": "수치 의미"
    }}
  ],
  "timeline": [
    {{
      "label": "날짜 또는 단계",
      "detail": "설명"
    }}
  ],
  "slide_plan": [
    {{
      "title": "슬라이드 제목",
      "layout": "summary",
      "bullets": ["내용 1", "내용 2"]
    }}
  ]
}}

slide_plan의 layout은
summary, issue_cause_action, timeline, chart, action_items, photo_text, conclusion
중 하나만 사용하세요.
"""


def analyze_media(
    uploaded_file,
    analysis_type="회의록",
    language_mode="한국어 + 영어 혼용",
    glossary=None,
):
    if uploaded_file is None:
        raise ValueError("분석할 음성 또는 영상 파일을 선택해주세요.")

    glossary = glossary or []
    suffix = Path(uploaded_file.name).suffix or ".bin"

    client = _client()
    temp_path = None
    remote_file = None

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(uploaded_file.getbuffer())
            temp_path = tmp.name

        remote_file = client.files.upload(file=temp_path)

        start_time = time.time()

        while True:
            state = getattr(remote_file, "state", None)
            state_name = getattr(state, "name", "") if state else ""

            if state_name in ("", "ACTIVE"):
                break

            if state_name == "FAILED":
                raise RuntimeError(
                    "Gemini가 업로드한 미디어 파일을 처리하지 못했습니다."
                )

            if time.time() - start_time > 600:
                raise TimeoutError(
                    "미디어 처리 시간이 10분을 초과했습니다."
                )

            time.sleep(3)
            remote_file = client.files.get(name=remote_file.name)

        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=[
                _media_prompt(
                    analysis_type,
                    language_mode,
                    glossary,
                ),
                remote_file,
            ],
            config={
                "response_mime_type": "application/json",
            },
        )

        return _normalize(_extract_json(response.text))

    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except OSError:
                pass

        if remote_file is not None:
            try:
                client.files.delete(name=remote_file.name)
            except Exception:
                pass


def rebuild_from_transcript(
    transcript,
    analysis_type="회의록",
    glossary=None,
):
    if not transcript.strip():
        raise ValueError("Transcript 내용이 없습니다.")

    glossary = glossary or []
    glossary_text = ", ".join(glossary) if glossary else "없음"

    prompt = f"""
다음 Transcript는 사용자가 직접 검토하고 수정한 최종 원문입니다.
이 Transcript만 사실 근거로 사용해서 요약, Action Item, 숫자, Timeline, PPT 구성안을 다시 만드세요.

[분석 유형]
{analysis_type}

[전문용어]
{glossary_text}

[Transcript]
{transcript}

없는 사실, 담당자, 날짜, 숫자를 만들지 마세요.

아래 JSON만 반환하세요.

{{
  "title": "분석 제목",
  "summary": ["핵심 요약"],
  "key_points": ["중요 포인트"],
  "action_items": [
    {{
      "owner": "담당 또는 미정",
      "task": "업무",
      "due": "기한 또는 미정",
      "status": "대기"
    }}
  ],
  "numbers": [
    {{
      "label": "항목",
      "value": 100,
      "unit": "단위",
      "context": "의미"
    }}
  ],
  "timeline": [
    {{
      "label": "날짜 또는 단계",
      "detail": "설명"
    }}
  ],
  "slide_plan": [
    {{
      "title": "슬라이드 제목",
      "layout": "summary",
      "bullets": ["내용"]
    }}
  ]
}}
"""

    client = _client()

    response = client.models.generate_content(
        model=MODEL_NAME,
        contents=prompt,
        config={
            "response_mime_type": "application/json",
        },
    )

    result = _normalize(_extract_json(response.text))
    result["transcript"] = transcript
    result["uncertain_terms"] = []

    return result


def _theme(theme_name):
    return THEMES.get(
        theme_name,
        THEMES["Corporate Basic"],
    )


def _set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=16,
    bold=False,
    color=RGBColor(45, 49, 58),
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    paragraph = frame.paragraphs[0]
    paragraph.alignment = align

    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    return box


def _add_header(slide, title, theme):
    _add_text(
        slide,
        title,
        0.72,
        0.42,
        11.8,
        0.55,
        24,
        True,
        theme["dark"],
    )

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.72),
        Inches(1.08),
        Inches(0.8),
        Inches(0.06),
    )

    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = theme["accent"]
    accent_bar.line.fill.background()


def _add_footer(slide, page_number):
    _add_text(
        slide,
        "AIAI Assistant",
        0.72,
        7.03,
        3.0,
        0.2,
        8,
        False,
        RGBColor(122, 128, 139),
    )

    _add_text(
        slide,
        f"Confidential · {page_number:02d}",
        10.2,
        7.03,
        2.4,
        0.2,
        8,
        False,
        RGBColor(122, 128, 139),
        PP_ALIGN.RIGHT,
    )


def _add_bullets(
    slide,
    items,
    x=0.9,
    y=1.55,
    w=11.2,
    h=4.9,
):
    items = [
        str(item)
        for item in items
        if str(item).strip()
    ]

    if not items:
        items = ["분석 결과가 없습니다."]

    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, item in enumerate(items[:7]):
        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )

        paragraph.text = f"• {item}"
        paragraph.font.name = "Malgun Gothic"
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(55, 60, 70)
        paragraph.space_after = Pt(10)


def _add_summary_cards(
    slide,
    summary,
    theme,
):
    items = list(summary[:3])

    while len(items) < 3:
        items.append("추가 핵심 내용 없음")

    positions = [
        0.78,
        4.48,
        8.18,
    ]

    for index, item in enumerate(items):
        x = positions[index]

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.65),
            Inches(3.35),
            Inches(3.8),
        )

        card.fill.solid()
        card.fill.fore_color.rgb = theme["soft"]
        card.line.color.rgb = RGBColor(225, 228, 235)

        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.22),
            Inches(1.92),
            Inches(0.55),
            Inches(0.4),
        )

        badge.fill.solid()
        badge.fill.fore_color.rgb = theme["accent"]
        badge.line.fill.background()

        _add_text(
            slide,
            f"{index + 1:02d}",
            x + 0.22,
            2.0,
            0.55,
            0.18,
            9,
            True,
            RGBColor(255, 255, 255),
            PP_ALIGN.CENTER,
        )

        _add_text(
            slide,
            f"Key Point {index + 1}",
            x + 0.92,
            1.93,
            2.1,
            0.35,
            13,
            True,
            theme["dark"],
        )

        _add_text(
            slide,
            item,
            x + 0.25,
            2.65,
            2.85,
            2.25,
            12,
            False,
            RGBColor(80, 86, 98),
        )


def _numeric_rows(numbers):
    rows = []

    for item in numbers:
        try:
            value = float(item.get("value"))
        except (TypeError, ValueError):
            continue

        label = (
            str(item.get("label", "")).strip()
            or "항목"
        )

        rows.append(
            {
                "label": label[:24],
                "value": value,
                "unit": str(item.get("unit", "")),
                "context": str(item.get("context", "")),
            }
        )

    return rows[:8]


def build_pptx(
    analysis,
    theme_name="Corporate Basic",
    reference_images=None,
):
    theme = _theme(theme_name)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    blank = presentation.slide_layouts[6]
    page = 1

    slide = presentation.slides.add_slide(blank)
    _set_background(
        slide,
        RGBColor(255, 255, 255),
    )

    _add_text(
        slide,
        analysis.get(
            "title",
            "회의·발표 분석",
        ),
        0.92,
        2.22,
        11.4,
        0.9,
        31,
        True,
        theme["dark"],
    )

    _add_text(
        slide,
        "AI 기반 회의·발표 요약 및 시각화",
        0.92,
        3.18,
        11.4,
        0.45,
        16,
        False,
        RGBColor(105, 111, 122),
    )

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.92),
        Inches(4.02),
        Inches(1.45),
        Inches(0.08),
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = theme["accent"]
    accent.line.fill.background()

    _add_footer(
        slide,
        page,
    )

    page += 1

    slide = presentation.slides.add_slide(blank)
    _set_background(
        slide,
        RGBColor(255, 255, 255),
    )

    _add_header(
        slide,
        "Executive Summary",
        theme,
    )

    _add_summary_cards(
        slide,
        analysis.get("summary", []),
        theme,
    )

    _add_footer(
        slide,
        page,
    )

    page += 1

    numbers = _numeric_rows(
        analysis.get("numbers", [])
    )

    if len(numbers) >= 2:
        slide = presentation.slides.add_slide(blank)

        _set_background(
            slide,
            RGBColor(255, 255, 255),
        )

        _add_header(
            slide,
            "Key Data Visualization",
            theme,
        )

        chart_data = ChartData()
        chart_data.categories = [
            item["label"]
            for item in numbers
        ]

        chart_data.add_series(
            "Value",
            [
                item["value"]
                for item in numbers
            ],
        )

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.95),
            Inches(1.6),
            Inches(7.25),
            Inches(4.75),
            chart_data,
        ).chart

        chart.has_legend = False
        chart.value_axis.has_major_gridlines = True
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.size = Pt(10)

        try:
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = (
                theme["accent"]
            )
        except Exception:
            pass

        context_lines = []

        for item in numbers[:5]:
            text = item["label"]

            if item["unit"]:
                text += f' ({item["unit"]})'

            if item["context"]:
                text += f': {item["context"]}'

            context_lines.append(text)

        _add_bullets(
            slide,
            context_lines,
            8.55,
            1.7,
            3.8,
            4.5,
        )

        _add_footer(
            slide,
            page,
        )

        page += 1

    timeline = analysis.get(
        "timeline",
        [],
    )

    if timeline:
        slide = presentation.slides.add_slide(blank)

        _set_background(
            slide,
            RGBColor(255, 255, 255),
        )

        _add_header(
            slide,
            "Timeline / Process",
            theme,
        )

        items = timeline[:5]
        gap = 11.6 / max(len(items), 1)

        for index, item in enumerate(items):
            x = 0.82 + index * gap

            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(x),
                Inches(2.15),
                Inches(0.34),
                Inches(0.34),
            )

            dot.fill.solid()
            dot.fill.fore_color.rgb = theme["accent"]
            dot.line.fill.background()

            if index < len(items) - 1:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x + 0.3),
                    Inches(2.30),
                    Inches(
                        max(
                            gap - 0.25,
                            0.2,
                        )
                    ),
                    Inches(0.04),
                )

                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(
                    210,
                    214,
                    221,
                )
                line.line.fill.background()

            _add_text(
                slide,
                item.get(
                    "label",
                    "",
                ),
                x - 0.05,
                2.7,
                max(
                    gap - 0.1,
                    1.2,
                ),
                0.5,
                12,
                True,
                theme["dark"],
            )

            _add_text(
                slide,
                item.get(
                    "detail",
                    "",
                ),
                x - 0.05,
                3.25,
                max(
                    gap - 0.15,
                    1.2,
                ),
                1.5,
                10,
                False,
                RGBColor(
                    92,
                    98,
                    108,
                ),
            )

        _add_footer(
            slide,
            page,
        )

        page += 1

    actions = analysis.get(
        "action_items",
        [],
    )

    if actions:
        slide = presentation.slides.add_slide(blank)

        _set_background(
            slide,
            RGBColor(255, 255, 255),
        )

        _add_header(
            slide,
            "Action Items",
            theme,
        )

        rows = [
            [
                "Owner",
                "Task",
                "Due",
                "Status",
            ]
        ]

        for item in actions[:7]:
            rows.append(
                [
                    str(
                        item.get(
                            "owner",
                            "미정",
                        )
                    ),
                    str(
                        item.get(
                            "task",
                            "",
                        )
                    ),
                    str(
                        item.get(
                            "due",
                            "미정",
                        )
                    ),
                    str(
                        item.get(
                            "status",
                            "대기",
                        )
                    ),
                ]
            )

        widths = [
            1.8,
            6.0,
            2.0,
            1.6,
        ]

        y = 1.55

        for row_index, row in enumerate(rows):
            x = 0.82

            for col_index, value in enumerate(row):
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x),
                    Inches(y),
                    Inches(
                        widths[col_index]
                    ),
                    Inches(0.64),
                )

                cell.fill.solid()

                cell.fill.fore_color.rgb = (
                    theme["accent"]
                    if row_index == 0
                    else RGBColor(
                        250,
                        251,
                        253,
                    )
                )

                cell.line.color.rgb = RGBColor(
                    225,
                    228,
                    235,
                )

                _add_text(
                    slide,
                    value,
                    x + 0.08,
                    y + 0.17,
                    widths[col_index] - 0.16,
                    0.28,
                    10,
                    row_index == 0,
                    (
                        RGBColor(
                            255,
                            255,
                            255,
                        )
                        if row_index == 0
                        else RGBColor(
                            58,
                            63,
                            72,
                        )
                    ),
                )

                x += widths[col_index]

            y += 0.64

        _add_footer(
            slide,
            page,
        )

        page += 1

    slide_plan = analysis.get(
        "slide_plan",
        [],
    )

    for plan in slide_plan[:4]:
        title = str(
            plan.get(
                "title",
                "",
            )
        ).strip()

        if not title:
            continue

        slide = presentation.slides.add_slide(blank)

        _set_background(
            slide,
            RGBColor(
                255,
                255,
                255,
            ),
        )

        _add_header(
            slide,
            title,
            theme,
        )

        _add_bullets(
            slide,
            plan.get(
                "bullets",
                [],
            ),
        )

        _add_footer(
            slide,
            page,
        )

        page += 1

    if reference_images:
        valid_images = reference_images[:4]

        if valid_images:
            slide = presentation.slides.add_slide(blank)

            _set_background(
                slide,
                RGBColor(
                    255,
                    255,
                    255,
                ),
            )

            _add_header(
                slide,
                "Reference Images",
                theme,
            )

            positions = [
                (0.9, 1.55),
                (6.78, 1.55),
                (0.9, 4.18),
                (6.78, 4.18),
            ]

            for image_item, position in zip(
                valid_images,
                positions,
            ):
                try:
                    image_bytes = (
                        image_item["bytes"]
                        if isinstance(
                            image_item,
                            dict,
                        )
                        else image_item
                    )

                    slide.shapes.add_picture(
                        BytesIO(image_bytes),
                        Inches(position[0]),
                        Inches(position[1]),
                        width=Inches(5.25),
                        height=Inches(2.15),
                    )

                except Exception:
                    pass

            _add_footer(
                slide,
                page,
            )

            page += 1

    slide = presentation.slides.add_slide(blank)

    _set_background(
        slide,
        RGBColor(
            255,
            255,
            255,
        ),
    )

    _add_header(
        slide,
        "Conclusion & Next Step",
        theme,
    )

    _add_bullets(
        slide,
        analysis.get(
            "key_points",
            analysis.get(
                "summary",
                [],
            ),
        ),
    )

    _add_footer(
        slide,
        page,
    )

    output = BytesIO()
    presentation.save(output)

    return output.getvalue()


def build_pdf(
    analysis,
    report_title=None,
):
    output = BytesIO()

    pdfmetrics.registerFont(
        UnicodeCIDFont(
            "HYSMyeongJo-Medium"
        )
    )

    document = SimpleDocTemplate(
        output,
        pagesize=landscape(A4),
        rightMargin=34,
        leftMargin=34,
        topMargin=32,
        bottomMargin=32,
    )

    styles = getSampleStyleSheet()

    body_style = ParagraphStyle(
        "KoreanBody",
        parent=styles["BodyText"],
        fontName="HYSMyeongJo-Medium",
        fontSize=10,
        leading=15,
        textColor=colors.HexColor(
            "#30343B"
        ),
    )

    title_style = ParagraphStyle(
        "KoreanTitle",
        parent=body_style,
        fontSize=22,
        leading=28,
        spaceAfter=14,
        textColor=colors.HexColor(
            "#1F2937"
        ),
    )

    heading_style = ParagraphStyle(
        "KoreanHeading",
        parent=body_style,
        fontSize=14,
        leading=19,
        spaceBefore=10,
        spaceAfter=8,
        textColor=colors.HexColor(
            "#4850B4"
        ),
    )

    story = []

    story.append(
        Paragraph(
            (
                report_title
                or analysis.get(
                    "title",
                    "회의·발표 분석",
                )
            ),
            title_style,
        )
    )

    story.append(
        Paragraph(
            "AI 기반 회의·발표 요약 및 검토 보고서",
            body_style,
        )
    )

    story.append(
        Spacer(
            1,
            14,
        )
    )

    story.append(
        Paragraph(
            "Executive Summary",
            heading_style,
        )
    )

    for item in analysis.get(
        "summary",
        [],
    ):
        story.append(
            Paragraph(
                f"• {str(item)}",
                body_style,
            )
        )

    actions = analysis.get(
        "action_items",
        [],
    )

    if actions:
        story.append(
            Spacer(
                1,
                12,
            )
        )

        story.append(
            Paragraph(
                "Action Items",
                heading_style,
            )
        )

        table_data = [
            [
                "담당",
                "업무",
                "기한",
                "상태",
            ]
        ]

        for item in actions[:12]:
            table_data.append(
                [
                    str(
                        item.get(
                            "owner",
                            "미정",
                        )
                    ),
                    str(
                        item.get(
                            "task",
                            "",
                        )
                    ),
                    str(
                        item.get(
                            "due",
                            "미정",
                        )
                    ),
                    str(
                        item.get(
                            "status",
                            "대기",
                        )
                    ),
                ]
            )

        table = Table(
            table_data,
            colWidths=[
                85,
                330,
                100,
                90,
            ],
            repeatRows=1,
        )

        table.setStyle(
            TableStyle(
                [
                    (
                        "FONTNAME",
                        (0, 0),
                        (-1, -1),
                        "HYSMyeongJo-Medium",
                    ),
                    (
                        "FONTSIZE",
                        (0, 0),
                        (-1, -1),
                        9,
                    ),
                    (
                        "BACKGROUND",
                        (0, 0),
                        (-1, 0),
                        colors.HexColor(
                            "#4850B4"
                        ),
                    ),
                    (
                        "TEXTCOLOR",
                        (0, 0),
                        (-1, 0),
                        colors.white,
                    ),
                    (
                        "GRID",
                        (0, 0),
                        (-1, -1),
                        0.5,
                        colors.HexColor(
                            "#D9DDE5"
                        ),
                    ),
                    (
                        "VALIGN",
                        (0, 0),
                        (-1, -1),
                        "TOP",
                    ),
                    (
                        "LEFTPADDING",
                        (0, 0),
                        (-1, -1),
                        7,
                    ),
                    (
                        "RIGHTPADDING",
                        (0, 0),
                        (-1, -1),
                        7,
                    ),
                    (
                        "TOPPADDING",
                        (0, 0),
                        (-1, -1),
                        7,
                    ),
                    (
                        "BOTTOMPADDING",
                        (0, 0),
                        (-1, -1),
                        7,
                    ),
                ]
            )
        )

        story.append(table)

    timeline = analysis.get(
        "timeline",
        [],
    )

    if timeline:
        story.append(
            Spacer(
                1,
                12,
            )
        )

        story.append(
            Paragraph(
                "Timeline / Process",
                heading_style,
            )
        )

        for item in timeline:
            story.append(
                Paragraph(
                    (
                        f'• {item.get("label", "")} '
                        f'— {item.get("detail", "")}'
                    ),
                    body_style,
                )
            )

    story.append(
        PageBreak()
    )

    story.append(
        Paragraph(
            "Transcript",
            title_style,
        )
    )

    transcript = str(
        analysis.get(
            "transcript",
            "",
        )
    ).replace(
        "\n",
        "<br/>",
    )

    story.append(
        Paragraph(
            (
                transcript
                if transcript
                else "Transcript 없음"
            ),
            body_style,
        )
    )

    document.build(story)

    return output.getvalue()
